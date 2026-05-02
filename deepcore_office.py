"""
deepcore_office.py — Motor de exportación Office con IA para DeepCore
Genera Excel profesional, reportes Word y presentaciones PowerPoint con análisis de IA.
Compatible con todos los programas del suite DeepCore.

Uso:
    from deepcore_office import DeepCoreOffice
    office = DeepCoreOffice(api_key="sk-ant-...")
    path = office.exportar_excel(titulo="Inventario", columnas=[...], filas=[...], kpis={...})
"""

import os
import threading
from datetime import datetime
from typing import Optional

# ── Rutas de salida ───────────────────────────────────────────────────────────
_EXPORTS_DIR = os.path.join(
    os.environ.get('USERPROFILE', os.path.expanduser('~')),
    'Documents', 'DeepCore', 'Exports'
)

# ── Paleta DeepCore (colores para Excel/PPT) ──────────────────────────────────
_NAVY      = '0F172A'
_SLATE     = '1E293B'
_BLUE      = '3B82F6'
_GREEN     = '22C55E'
_TEAL      = '14B8A6'
_AMBER     = 'F59E0B'
_RED       = 'EF4444'
_TEXT      = 'F8FAFC'
_SUBTEXT   = '94A3B8'
_BORDER    = '334155'
_WHITE     = 'FFFFFF'
_DARK_ROW  = '0F172A'
_ALT_ROW   = '1E293B'


def _ensure_dir(path: str) -> str:
    os.makedirs(path, exist_ok=True)
    return path


def _ts() -> str:
    return datetime.now().strftime('%Y%m%d_%H%M%S')


# ── Análisis con IA (Claude claude-haiku-4-5) ──────────────────────────────────────
def _ai_analizar(prompt: str, api_key: str, max_tokens: int = 600) -> str:
    if not api_key:
        return ''
    try:
        import anthropic
        client = anthropic.Anthropic(api_key=api_key)
        resp = client.messages.create(
            model='claude-haiku-4-5-20251001',
            max_tokens=max_tokens,
            messages=[{'role': 'user', 'content': prompt}]
        )
        return resp.content[0].text.strip()
    except Exception as e:
        return f'(Análisis IA no disponible: {e})'


# ═══════════════════════════════════════════════════════════════════════════════
# EXCEL
# ═══════════════════════════════════════════════════════════════════════════════

def _excel_color(hex_str: str):
    from openpyxl.styles import PatternFill
    return PatternFill(start_color=hex_str, end_color=hex_str, fill_type='solid')


def _excel_font(bold=False, size=11, color=_WHITE, name='Segoe UI'):
    from openpyxl.styles import Font
    return Font(name=name, size=size, bold=bold, color=color)


def _excel_align(h='left', v='center', wrap=False):
    from openpyxl.styles import Alignment
    return Alignment(horizontal=h, vertical=v, wrap_text=wrap)


def _excel_border():
    from openpyxl.styles import Border, Side
    thin = Side(border_style='thin', color=_BORDER)
    return Border(bottom=thin)


def exportar_excel(
    titulo: str,
    columnas: list,
    filas: list,
    kpis: Optional[dict] = None,
    programa: str = 'DeepCore',
    api_key: str = '',
    subtitulo: str = '',
) -> str:
    """
    Genera un .xlsx profesional con marca DeepCore, datos formateados y análisis IA.

    Args:
        titulo: Nombre del reporte (ej: 'Reporte de Inventario')
        columnas: Lista de strings con los encabezados
        filas: Lista de listas/tuplas con los datos
        kpis: Dict con métricas clave {label: valor}
        programa: Nombre del programa DeepCore
        api_key: Clave Anthropic para análisis IA (opcional)
        subtitulo: Texto adicional bajo el título

    Returns:
        Ruta absoluta del archivo generado
    """
    from openpyxl import Workbook
    from openpyxl.styles import PatternFill, Font, Alignment, Border, Side
    from openpyxl.utils import get_column_letter
    from openpyxl.chart import BarChart, Reference

    _ensure_dir(_EXPORTS_DIR)
    fecha = datetime.now().strftime('%d/%m/%Y %H:%M')
    nombre_archivo = f"{titulo.replace(' ', '_')}_{_ts()}.xlsx"
    ruta = os.path.join(_EXPORTS_DIR, nombre_archivo)

    wb = Workbook()

    # ── Hoja 1: Datos ──────────────────────────────────────────────────────────
    ws = wb.active
    ws.title = 'Datos'
    ws.sheet_view.showGridLines = False

    # Fila 1: Marca DeepCore
    ws.merge_cells(f'A1:{get_column_letter(max(len(columnas), 4))}1')
    c = ws['A1']
    c.value = f'DeepCore — {programa}'
    c.fill = _excel_color(_NAVY)
    c.font = Font(name='Segoe UI', size=9, bold=False, color=_SUBTEXT)
    c.alignment = Alignment(horizontal='left', vertical='center')
    ws.row_dimensions[1].height = 18

    # Fila 2: Título principal
    ws.merge_cells(f'A2:{get_column_letter(max(len(columnas), 4))}2')
    c = ws['A2']
    c.value = titulo.upper()
    c.fill = _excel_color(_NAVY)
    c.font = Font(name='Segoe UI', size=18, bold=True, color=_TEXT)
    c.alignment = Alignment(horizontal='left', vertical='center')
    ws.row_dimensions[2].height = 36

    # Fila 3: Subtítulo + fecha
    col_total = max(len(columnas), 4)
    mid = col_total // 2
    ws.merge_cells(f'A3:{get_column_letter(mid)}3')
    c = ws['A3']
    c.value = subtitulo or f'Generado el {fecha}'
    c.fill = _excel_color(_SLATE)
    c.font = Font(name='Segoe UI', size=10, color=_SUBTEXT)
    c.alignment = Alignment(horizontal='left', vertical='center')

    ws.merge_cells(f'{get_column_letter(mid+1)}3:{get_column_letter(col_total)}3')
    c = ws.cell(row=3, column=mid+1)
    c.value = fecha
    c.fill = _excel_color(_SLATE)
    c.font = Font(name='Segoe UI', size=10, color=_SUBTEXT)
    c.alignment = Alignment(horizontal='right', vertical='center')
    ws.row_dimensions[3].height = 22

    # Fila 4: separador
    for col in range(1, col_total + 1):
        c = ws.cell(row=4, column=col)
        c.fill = _excel_color(_BLUE)
    ws.row_dimensions[4].height = 4

    # KPIs — dashboard profesional sin espacios muertos
    kpi_row_end = 4
    kpi_colors = [_GREEN, _BLUE, _AMBER, _TEAL, _RED, _GREEN, _BLUE, _AMBER]
    if kpis:
        kpi_keys = list(kpis.keys())
        kpi_vals = [str(kpis[k]) for k in kpi_keys]
        n_kpis = min(len(kpi_keys), col_total)

        # Fila 5: barras de acento por KPI (color identificador)
        for col in range(1, col_total + 1):
            col_hex = kpi_colors[(col - 1) % len(kpi_colors)] if col <= n_kpis else _SLATE
            ws.cell(row=5, column=col).fill = _excel_color(col_hex)
        ws.row_dimensions[5].height = 4

        # Fila 6: etiquetas — texto blanco claro y legible
        for i, label in enumerate(kpi_keys[:n_kpis], 1):
            c = ws.cell(row=6, column=i)
            c.value = label.upper()
            c.fill = _excel_color(_SLATE)
            c.font = Font(name='Segoe UI', size=11, bold=True, color=_TEXT)
            c.alignment = Alignment(horizontal='center', vertical='center', wrap_text=True)
        for col in range(n_kpis + 1, col_total + 1):
            ws.cell(row=6, column=col).fill = _excel_color(_SLATE)
        ws.row_dimensions[6].height = 30

        # Fila 7: valores grandes en color de acento
        for i, (val, col_hex) in enumerate(zip(kpi_vals[:n_kpis], kpi_colors), 1):
            c = ws.cell(row=7, column=i)
            c.value = val
            c.fill = _excel_color(_NAVY)
            c.font = Font(name='Segoe UI', size=20, bold=True, color=col_hex)
            c.alignment = Alignment(horizontal='center', vertical='center')
        for col in range(n_kpis + 1, col_total + 1):
            ws.cell(row=7, column=col).fill = _excel_color(_NAVY)
        ws.row_dimensions[7].height = 38

        # Fila 8: línea separadora
        for col in range(1, col_total + 1):
            ws.cell(row=8, column=col).fill = _excel_color(_BORDER)
        ws.row_dimensions[8].height = 3

        kpi_row_end = 8

    data_start = kpi_row_end + 1

    # Encabezados de columna
    col_colors = [_BLUE, _SLATE, _SLATE]
    for i, col_name in enumerate(columnas, 1):
        c = ws.cell(row=data_start, column=i)
        c.value = col_name.upper()
        c.fill = _excel_color(_BLUE if i == 1 else _SLATE)
        c.font = Font(name='Segoe UI', size=11, bold=True, color=_WHITE)
        c.alignment = Alignment(horizontal='left', vertical='center', wrap_text=False)
    ws.row_dimensions[data_start].height = 28

    # Filas de datos
    for row_idx, fila in enumerate(filas):
        r = data_start + 1 + row_idx
        bg = _DARK_ROW if row_idx % 2 == 0 else _ALT_ROW
        for col_idx, val in enumerate(fila, 1):
            c = ws.cell(row=r, column=col_idx)
            c.value = val
            c.fill = _excel_color(bg)
            c.font = Font(name='Segoe UI', size=11, color=_TEXT)
            c.alignment = Alignment(horizontal='left', vertical='center')
        ws.row_dimensions[r].height = 22

    # Autoajustar anchos
    for i, col_name in enumerate(columnas, 1):
        max_len = len(str(col_name))
        for fila in filas:
            if i <= len(fila):
                max_len = max(max_len, len(str(fila[i-1] or '')))
        ws.column_dimensions[get_column_letter(i)].width = min(max_len + 4, 40)

    # Congelar encabezados
    ws.freeze_panes = ws.cell(row=data_start + 1, column=1)

    # ── Hoja 2: Análisis IA ────────────────────────────────────────────────────
    if api_key:
        ws_ia = wb.create_sheet('Analisis IA')
        ws_ia.sheet_view.showGridLines = False

        # Construir prompt
        kpi_txt = '\n'.join([f'- {k}: {v}' for k, v in (kpis or {}).items()])
        filas_txt = '\n'.join([', '.join(str(v) for v in f) for f in filas[:20]])
        prompt = (
            f"Eres un analista de negocios experto para PyMEs ecuatorianas.\n"
            f"Analiza los siguientes datos del programa DeepCore {programa}:\n\n"
            f"Reporte: {titulo}\n"
            f"Fecha: {fecha}\n\n"
            f"INDICADORES CLAVE:\n{kpi_txt}\n\n"
            f"MUESTRA DE DATOS (primeras 20 filas):\n"
            f"Columnas: {', '.join(columnas)}\n"
            f"{filas_txt}\n\n"
            f"Proporciona:\n"
            f"1. RESUMEN EJECUTIVO (2-3 oraciones)\n"
            f"2. 3 HALLAZGOS PRINCIPALES\n"
            f"3. 3 RECOMENDACIONES ACCIONABLES para el negocio\n"
            f"4. ALERTAS o riesgos detectados (si los hay)\n\n"
            f"Responde en español, de forma concisa y práctica para un empresario."
        )
        analisis = _ai_analizar(prompt, api_key, max_tokens=800)

        # Encabezado hoja IA
        ws_ia.merge_cells('A1:G1')
        c = ws_ia['A1']
        c.value = 'ANALISIS INTELIGENTE — DEEPCORE IA'
        c.fill = _excel_color(_NAVY)
        c.font = Font(name='Segoe UI', size=16, bold=True, color=_BLUE)
        c.alignment = Alignment(horizontal='left', vertical='center')
        ws_ia.row_dimensions[1].height = 36

        ws_ia.merge_cells('A2:G2')
        c = ws_ia['A2']
        c.value = f'Generado automáticamente para: {titulo} | {fecha}'
        c.fill = _excel_color(_SLATE)
        c.font = Font(name='Segoe UI', size=10, color=_SUBTEXT)
        c.alignment = Alignment(horizontal='left', vertical='center')
        ws_ia.row_dimensions[2].height = 20

        # Contenido del análisis
        lineas = analisis.split('\n')
        row_ia = 4
        for linea in lineas:
            ws_ia.merge_cells(f'A{row_ia}:G{row_ia}')
            c = ws_ia.cell(row=row_ia, column=1)
            c.value = linea
            c.fill = _excel_color(_DARK_ROW)
            is_header = linea.strip().endswith(':') or (linea.strip() and linea.strip()[0].isdigit() and '.' in linea[:3])
            c.font = Font(
                name='Segoe UI',
                size=11 if is_header else 10,
                bold=is_header,
                color=_BLUE if is_header else _TEXT
            )
            c.alignment = Alignment(horizontal='left', vertical='center', wrap_text=True)
            ws_ia.row_dimensions[row_ia].height = 20 if not is_header else 24
            row_ia += 1

        ws_ia.column_dimensions['A'].width = 80

        # Nota de pie
        ws_ia.merge_cells(f'A{row_ia+2}:G{row_ia+2}')
        c = ws_ia.cell(row=row_ia+2, column=1)
        c.value = '* Análisis generado por DeepCore IA (Claude). Úsalo como guía, no como asesoramiento definitivo.'
        c.fill = _excel_color(_NAVY)
        c.font = Font(name='Segoe UI', size=9, italic=True, color=_SUBTEXT)
        c.alignment = Alignment(horizontal='left', vertical='center')

    wb.save(ruta)
    return ruta


# ═══════════════════════════════════════════════════════════════════════════════
# WORD
# ═══════════════════════════════════════════════════════════════════════════════

def exportar_word(
    titulo: str,
    columnas: list,
    filas: list,
    kpis: Optional[dict] = None,
    programa: str = 'DeepCore',
    api_key: str = '',
    subtitulo: str = '',
) -> str:
    """
    Genera un .docx profesional con marca DeepCore, tablas de datos y resumen IA.
    """
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT, WD_ALIGN_VERTICAL
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    import struct

    _ensure_dir(_EXPORTS_DIR)
    fecha = datetime.now().strftime('%d/%m/%Y %H:%M')
    nombre_archivo = f"{titulo.replace(' ', '_')}_{_ts()}.docx"
    ruta = os.path.join(_EXPORTS_DIR, nombre_archivo)

    doc = Document()

    # Márgenes
    for section in doc.sections:
        section.top_margin    = Cm(2)
        section.bottom_margin = Cm(2)
        section.left_margin   = Cm(2.5)
        section.right_margin  = Cm(2.5)

    def _set_color(run, hex_str: str):
        r, g, b = int(hex_str[0:2],16), int(hex_str[2:4],16), int(hex_str[4:6],16)
        run.font.color.rgb = RGBColor(r, g, b)

    def _set_bg(paragraph, hex_str: str):
        pPr = paragraph._p.get_or_add_pPr()
        shd = OxmlElement('w:shd')
        shd.set(qn('w:val'), 'clear')
        shd.set(qn('w:color'), 'auto')
        shd.set(qn('w:fill'), hex_str)
        pPr.append(shd)

    def _set_cell_bg(cell, hex_str: str):
        tcPr = cell._tc.get_or_add_tcPr()
        shd = OxmlElement('w:shd')
        shd.set(qn('w:val'), 'clear')
        shd.set(qn('w:color'), 'auto')
        shd.set(qn('w:fill'), hex_str)
        tcPr.append(shd)

    # ── Encabezado de marca ────────────────────────────────────────────────────
    p = doc.add_paragraph()
    _set_bg(p, _NAVY)
    p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    run = p.add_run(f'DeepCore — {programa}')
    run.font.name = 'Segoe UI'
    run.font.size = Pt(9)
    _set_color(run, _SUBTEXT)
    p.paragraph_format.space_before = Pt(0)
    p.paragraph_format.space_after  = Pt(0)

    # ── Título ────────────────────────────────────────────────────────────────
    p = doc.add_paragraph()
    _set_bg(p, _NAVY)
    p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    run = p.add_run(titulo.upper())
    run.font.name = 'Segoe UI'
    run.font.size = Pt(22)
    run.font.bold = True
    _set_color(run, _TEXT)
    p.paragraph_format.space_before = Pt(0)
    p.paragraph_format.space_after  = Pt(0)

    # ── Subtítulo + fecha ─────────────────────────────────────────────────────
    p = doc.add_paragraph()
    _set_bg(p, _SLATE)
    p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    run = p.add_run(f'{subtitulo or "Reporte generado por DeepCore IA"}    |    {fecha}')
    run.font.name = 'Segoe UI'
    run.font.size = Pt(9)
    _set_color(run, _SUBTEXT)
    p.paragraph_format.space_before = Pt(0)
    p.paragraph_format.space_after  = Pt(8)

    # ── KPIs ──────────────────────────────────────────────────────────────────
    if kpis:
        p = doc.add_paragraph()
        run = p.add_run('INDICADORES CLAVE')
        run.font.name = 'Segoe UI'
        run.font.size = Pt(9)
        run.font.bold = True
        _set_color(run, _SUBTEXT)
        p.paragraph_format.space_after = Pt(4)

        n_kpis = len(kpis)
        tbl_kpi = doc.add_table(rows=2, cols=min(n_kpis, 5))
        tbl_kpi.alignment = WD_TABLE_ALIGNMENT.LEFT
        kpi_list = list(kpis.items())[:5]
        kpi_colors_hex = [_GREEN, _BLUE, _AMBER, _TEAL, _RED]
        for i, (label, val) in enumerate(kpi_list):
            lbl_cell = tbl_kpi.rows[0].cells[i]
            _set_cell_bg(lbl_cell, _SLATE)
            p = lbl_cell.paragraphs[0]
            run = p.add_run(label.upper())
            run.font.name = 'Segoe UI'; run.font.size = Pt(8); run.font.bold = True
            _set_color(run, _SUBTEXT)

            val_cell = tbl_kpi.rows[1].cells[i]
            _set_cell_bg(val_cell, _NAVY)
            p = val_cell.paragraphs[0]
            run = p.add_run(str(val))
            run.font.name = 'Segoe UI'; run.font.size = Pt(16); run.font.bold = True
            _set_color(run, kpi_colors_hex[i % len(kpi_colors_hex)])

        doc.add_paragraph().paragraph_format.space_after = Pt(6)

    # ── Análisis IA ───────────────────────────────────────────────────────────
    if api_key:
        kpi_txt = '\n'.join([f'- {k}: {v}' for k, v in (kpis or {}).items()])
        filas_txt = '\n'.join([', '.join(str(v) for v in f) for f in filas[:15]])
        prompt = (
            f"Eres consultor de negocios para PyMEs ecuatorianas.\n"
            f"Escribe un RESUMEN EJECUTIVO profesional (máx 200 palabras) para este reporte:\n\n"
            f"Programa: {programa}\nReporte: {titulo}\nFecha: {fecha}\n"
            f"KPIs: {kpi_txt}\n"
            f"Datos (primeras filas): {', '.join(columnas)}\n{filas_txt}\n\n"
            f"Incluye: situación actual, puntos destacados, y 2 recomendaciones concretas."
        )
        resumen_ia = _ai_analizar(prompt, api_key, max_tokens=400)

        p = doc.add_paragraph()
        run = p.add_run('RESUMEN EJECUTIVO — DEEPCORE IA')
        run.font.name = 'Segoe UI'; run.font.size = Pt(9); run.font.bold = True
        _set_color(run, _BLUE)
        p.paragraph_format.space_after = Pt(4)

        p = doc.add_paragraph()
        run = p.add_run(resumen_ia)
        run.font.name = 'Segoe UI'; run.font.size = Pt(10)
        _set_color(run, '1E293B')
        p.paragraph_format.space_after = Pt(10)

    # ── Tabla de datos ────────────────────────────────────────────────────────
    p = doc.add_paragraph()
    run = p.add_run('DATOS DEL REPORTE')
    run.font.name = 'Segoe UI'; run.font.size = Pt(9); run.font.bold = True
    _set_color(run, _SUBTEXT)
    p.paragraph_format.space_after = Pt(4)

    if columnas and filas:
        tbl = doc.add_table(rows=1 + len(filas), cols=len(columnas))
        tbl.alignment = WD_TABLE_ALIGNMENT.LEFT
        tbl.style = 'Table Grid'

        # Encabezado
        for i, col_name in enumerate(columnas):
            cell = tbl.rows[0].cells[i]
            _set_cell_bg(cell, _BLUE)
            p = cell.paragraphs[0]
            run = p.add_run(col_name.upper())
            run.font.name = 'Segoe UI'; run.font.size = Pt(9); run.font.bold = True
            _set_color(run, _WHITE)

        # Datos
        for row_i, fila in enumerate(filas):
            bg = _DARK_ROW if row_i % 2 == 0 else _ALT_ROW
            for col_i, val in enumerate(fila):
                if col_i >= len(columnas): break
                cell = tbl.rows[row_i+1].cells[col_i]
                _set_cell_bg(cell, bg)
                p = cell.paragraphs[0]
                run = p.add_run(str(val or ''))
                run.font.name = 'Segoe UI'; run.font.size = Pt(9)
                _set_color(run, _TEXT)

    # Pie de página
    doc.add_paragraph()
    p = doc.add_paragraph()
    run = p.add_run(f'DeepCore {programa} — {fecha} — deepcore.ec')
    run.font.name = 'Segoe UI'; run.font.size = Pt(8); run.font.italic = True
    _set_color(run, _SUBTEXT)
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.save(ruta)
    return ruta


# ═══════════════════════════════════════════════════════════════════════════════
# POWERPOINT
# ═══════════════════════════════════════════════════════════════════════════════

def exportar_pptx(
    titulo: str,
    columnas: list,
    filas: list,
    kpis: Optional[dict] = None,
    programa: str = 'DeepCore',
    api_key: str = '',
    subtitulo: str = '',
) -> str:
    """
    Genera una presentación .pptx de alta calidad con tema DeepCore dark.
    Incluye slide de título, KPIs, tabla de datos y recomendaciones IA.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    def _rgb(hex_str):
        return RGBColor(int(hex_str[0:2],16), int(hex_str[2:4],16), int(hex_str[4:6],16))

    def _set_cell_bg(cell, hex_str):
        from pptx.oxml.ns import qn
        from lxml import etree
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', hex_str)

    def _add_text_box(slide, text, left, top, width, height,
                      font_size=18, bold=False, color=_TEXT, align=PP_ALIGN.LEFT,
                      bg_color=None):
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height))
        if bg_color:
            fill = txBox.fill
            fill.solid()
            fill.fore_color.rgb = _rgb(bg_color)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = _rgb(color)
        run.font.name = 'Segoe UI'
        return txBox

    _ensure_dir(_EXPORTS_DIR)
    fecha = datetime.now().strftime('%d/%m/%Y')
    nombre_archivo = f"{titulo.replace(' ', '_')}_{_ts()}.pptx"
    ruta = os.path.join(_EXPORTS_DIR, nombre_archivo)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # Completamente en blanco

    def _bg(slide, hex_str):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = _rgb(hex_str)

    # ── Slide 1: Portada ───────────────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank_layout)
    _bg(s1, _NAVY)

    # Barra lateral izquierda de color
    bar = s1.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0), Inches(0.15), Inches(7.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = _rgb(_BLUE)
    bar.line.fill.background()

    # Logo/badge DeepCore
    badge = s1.shapes.add_shape(1, Inches(0.5), Inches(0.4), Inches(1.4), Inches(0.5))
    badge.fill.solid(); badge.fill.fore_color.rgb = _rgb(_BLUE)
    badge.line.fill.background()
    tf = badge.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = 'DeepCore'; run.font.size = Pt(11)
    run.font.bold = True; run.font.color.rgb = _rgb(_NAVY); run.font.name = 'Segoe UI'

    _add_text_box(s1, programa.upper(), 0.5, 1.2, 8, 0.5,
                  font_size=12, color=_BLUE, bold=True)
    _add_text_box(s1, titulo, 0.5, 1.8, 10, 1.5,
                  font_size=36, bold=True, color=_TEXT)
    _add_text_box(s1, subtitulo or f'Reporte de negocios — {fecha}',
                  0.5, 3.5, 8, 0.5, font_size=14, color=_SUBTEXT)
    _add_text_box(s1, fecha, 0.5, 6.8, 5, 0.4, font_size=10, color=_SUBTEXT)

    # Círculo decorativo
    circ = s1.shapes.add_shape(9, Inches(10.5), Inches(1.5), Inches(2.5), Inches(2.5))
    circ.fill.solid(); circ.fill.fore_color.rgb = _rgb(_SLATE)
    circ.line.color.rgb = _rgb(_BLUE); circ.line.width = Pt(1)

    # ── Slide 2: KPIs / Dashboard ──────────────────────────────────────────────
    if kpis:
        s2 = prs.slides.add_slide(blank_layout)
        _bg(s2, _NAVY)
        bar2 = s2.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.08))
        bar2.fill.solid(); bar2.fill.fore_color.rgb = _rgb(_BLUE); bar2.line.fill.background()

        _add_text_box(s2, 'INDICADORES CLAVE', 0.4, 0.2, 6, 0.5,
                      font_size=10, color=_SUBTEXT, bold=True)
        _add_text_box(s2, titulo, 0.4, 0.55, 8, 0.7,
                      font_size=22, bold=True, color=_TEXT)

        kpi_list = list(kpis.items())
        kpi_colors_hex = [_GREEN, _BLUE, _AMBER, _TEAL, _RED, _GREEN, _BLUE, _AMBER]
        cols = min(len(kpi_list), 4)
        card_w = 12.0 / cols
        card_h = 2.0
        for i, (label, val) in enumerate(kpi_list[:cols*2]):
            row = i // cols
            col = i % cols
            x = 0.4 + col * card_w
            y = 1.5 + row * (card_h + 0.2)
            col_hex = kpi_colors_hex[i % len(kpi_colors_hex)]
            # Card background
            card_bg = s2.shapes.add_shape(1, Inches(x), Inches(y), Inches(card_w-0.2), Inches(card_h))
            card_bg.fill.solid(); card_bg.fill.fore_color.rgb = _rgb(_SLATE)
            card_bg.line.color.rgb = _rgb(_BORDER); card_bg.line.width = Pt(0.5)
            # Barra de color superior
            card_bar = s2.shapes.add_shape(1, Inches(x), Inches(y), Inches(card_w-0.2), Inches(0.05))
            card_bar.fill.solid(); card_bar.fill.fore_color.rgb = _rgb(col_hex)
            card_bar.line.fill.background()
            # Textos
            _add_text_box(s2, label.upper(), x+0.15, y+0.15, card_w-0.4, 0.35,
                          font_size=9, color=_SUBTEXT, bold=True)
            _add_text_box(s2, str(val), x+0.15, y+0.55, card_w-0.4, 0.9,
                          font_size=26, bold=True, color=col_hex)

        _add_text_box(s2, f'DeepCore {programa}  |  {fecha}', 0.4, 7.1, 8, 0.3,
                      font_size=8, color=_SUBTEXT)

    # ── Slide 3: Tabla de datos ────────────────────────────────────────────────
    if columnas and filas:
        s3 = prs.slides.add_slide(blank_layout)
        _bg(s3, _NAVY)
        bar3 = s3.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.08))
        bar3.fill.solid(); bar3.fill.fore_color.rgb = _rgb(_BLUE); bar3.line.fill.background()

        _add_text_box(s3, 'DATOS DEL REPORTE', 0.4, 0.2, 6, 0.4,
                      font_size=10, color=_SUBTEXT, bold=True)
        _add_text_box(s3, titulo, 0.4, 0.55, 10, 0.55,
                      font_size=18, bold=True, color=_TEXT)

        n_cols = len(columnas)
        n_rows = min(len(filas), 14)  # máx 14 filas en slide
        tbl_w = 12.4
        tbl_h = 5.5
        col_w = tbl_w / n_cols

        tbl_shape = s3.shapes.add_table(n_rows + 1, n_cols,
                                         Inches(0.4), Inches(1.4),
                                         Inches(tbl_w), Inches(tbl_h))
        tbl = tbl_shape.table

        for c_i, col_name in enumerate(columnas):
            cell = tbl.cell(0, c_i)
            _set_cell_bg(cell, _BLUE)
            tf = cell.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            run = p.add_run(); run.text = col_name.upper()
            run.font.bold = True; run.font.size = Pt(9)
            run.font.color.rgb = _rgb(_WHITE); run.font.name = 'Segoe UI'

        for r_i, fila in enumerate(filas[:n_rows]):
            bg = _DARK_ROW if r_i % 2 == 0 else _ALT_ROW
            for c_i, val in enumerate(fila[:n_cols]):
                cell = tbl.cell(r_i + 1, c_i)
                _set_cell_bg(cell, bg)
                tf = cell.text_frame; p = tf.paragraphs[0]
                run = p.add_run(); run.text = str(val or '')
                run.font.size = Pt(8); run.font.color.rgb = _rgb(_TEXT)
                run.font.name = 'Segoe UI'

        if len(filas) > n_rows:
            _add_text_box(s3, f'... y {len(filas)-n_rows} registros más',
                          0.4, 7.0, 8, 0.3, font_size=9, color=_SUBTEXT)

        _add_text_box(s3, f'DeepCore {programa}  |  {fecha}', 0.4, 7.1, 8, 0.3,
                      font_size=8, color=_SUBTEXT)

    # ── Slide 4: Recomendaciones IA ────────────────────────────────────────────
    if api_key:
        kpi_txt = '\n'.join([f'{k}: {v}' for k, v in (kpis or {}).items()])
        prompt = (
            f"Eres consultor estratégico para PyMEs ecuatorianas.\n"
            f"Con base en este reporte de {programa} — {titulo}:\n"
            f"KPIs: {kpi_txt}\n"
            f"Total de registros: {len(filas)}\n\n"
            f"Proporciona EXACTAMENTE:\n"
            f"TITULO: (3-5 palabras impactantes)\n"
            f"RECOMENDACION 1: (acción concreta)\n"
            f"RECOMENDACION 2: (acción concreta)\n"
            f"RECOMENDACION 3: (acción concreta)\n"
            f"OPORTUNIDAD: (una oportunidad de crecimiento específica)\n\n"
            f"Máximo 15 palabras por punto. Español. Sin explicaciones extra."
        )
        ia_resp = _ai_analizar(prompt, api_key, max_tokens=300)

        s4 = prs.slides.add_slide(blank_layout)
        _bg(s4, _NAVY)
        bar4 = s4.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.08))
        bar4.fill.solid(); bar4.fill.fore_color.rgb = _rgb(_BLUE); bar4.line.fill.background()

        _add_text_box(s4, 'DEEPCORE IA', 0.4, 0.2, 6, 0.4,
                      font_size=10, color=_BLUE, bold=True)
        _add_text_box(s4, 'Análisis y Recomendaciones', 0.4, 0.55, 10, 0.65,
                      font_size=22, bold=True, color=_TEXT)

        # Parsear respuesta IA
        lineas_ia = [l.strip() for l in ia_resp.split('\n') if l.strip()]
        rec_colors = [_GREEN, _TEAL, _BLUE, _AMBER]
        y_pos = 1.5
        for i, linea in enumerate(lineas_ia[:6]):
            is_title = linea.upper().startswith('TITULO')
            is_op    = linea.upper().startswith('OPORTUNIDAD')
            col_hex  = _BLUE if is_title else (_AMBER if is_op else rec_colors[i % len(rec_colors)])

            card = s4.shapes.add_shape(1, Inches(0.4), Inches(y_pos), Inches(12.5), Inches(0.75))
            card.fill.solid(); card.fill.fore_color.rgb = _rgb(_SLATE)
            card.line.fill.background()
            accent_bar = s4.shapes.add_shape(1, Inches(0.4), Inches(y_pos), Inches(0.06), Inches(0.75))
            accent_bar.fill.solid(); accent_bar.fill.fore_color.rgb = _rgb(col_hex)
            accent_bar.line.fill.background()

            _add_text_box(s4, linea, 0.7, y_pos+0.1, 12, 0.6,
                          font_size=12, color=_TEXT if not is_title else col_hex,
                          bold=is_title or is_op)
            y_pos += 0.9
            if y_pos > 6.5: break

        _add_text_box(s4, f'DeepCore {programa}  |  {fecha}  |  deepcore.ec',
                      0.4, 7.1, 10, 0.3, font_size=8, color=_SUBTEXT)

    # ── Slide 5: Cierre ────────────────────────────────────────────────────────
    s5 = prs.slides.add_slide(blank_layout)
    _bg(s5, _NAVY)
    circ_big = s5.shapes.add_shape(9, Inches(4.5), Inches(1.8), Inches(4.3), Inches(4.3))
    circ_big.fill.solid(); circ_big.fill.fore_color.rgb = _rgb(_SLATE)
    circ_big.line.color.rgb = _rgb(_BLUE); circ_big.line.width = Pt(1.5)
    _add_text_box(s5, 'DeepCore', 5.0, 3.0, 3.3, 0.7,
                  font_size=20, bold=True, color=_BLUE, align=PP_ALIGN.CENTER)
    _add_text_box(s5, programa, 5.0, 3.7, 3.3, 0.5,
                  font_size=13, color=_TEXT, align=PP_ALIGN.CENTER)
    _add_text_box(s5, 'deepcore.ec', 5.0, 4.3, 3.3, 0.4,
                  font_size=10, color=_SUBTEXT, align=PP_ALIGN.CENTER)
    _add_text_box(s5, f'Licencia $5/mes — Todo el suite', 3.5, 6.5, 6.3, 0.4,
                  font_size=9, color=_SUBTEXT, align=PP_ALIGN.CENTER)

    prs.save(ruta)
    return ruta


# ═══════════════════════════════════════════════════════════════════════════════
# CLASE PRINCIPAL
# ═══════════════════════════════════════════════════════════════════════════════

class DeepCoreOffice:
    """Fachada para exportar desde cualquier programa DeepCore."""

    def __init__(self, api_key: str = '', programa: str = 'DeepCore'):
        self.api_key  = api_key
        self.programa = programa

    def excel(self, titulo, columnas, filas, kpis=None, subtitulo='',
              callback=None):
        """Exporta a Excel en hilo separado. Llama callback(ruta) al terminar."""
        def _run():
            try:
                ruta = exportar_excel(titulo, columnas, filas, kpis=kpis,
                                      programa=self.programa, api_key=self.api_key,
                                      subtitulo=subtitulo)
                if callback: callback(ruta, None)
            except Exception as e:
                if callback: callback(None, str(e))
        threading.Thread(target=_run, daemon=True).start()

    def word(self, titulo, columnas, filas, kpis=None, subtitulo='',
             callback=None):
        """Exporta a Word en hilo separado."""
        def _run():
            try:
                ruta = exportar_word(titulo, columnas, filas, kpis=kpis,
                                     programa=self.programa, api_key=self.api_key,
                                     subtitulo=subtitulo)
                if callback: callback(ruta, None)
            except Exception as e:
                if callback: callback(None, str(e))
        threading.Thread(target=_run, daemon=True).start()

    def pptx(self, titulo, columnas, filas, kpis=None, subtitulo='',
             callback=None):
        """Crea presentación PowerPoint en hilo separado."""
        def _run():
            try:
                ruta = exportar_pptx(titulo, columnas, filas, kpis=kpis,
                                     programa=self.programa, api_key=self.api_key,
                                     subtitulo=subtitulo)
                if callback: callback(ruta, None)
            except Exception as e:
                if callback: callback(None, str(e))
        threading.Thread(target=_run, daemon=True).start()
